import json
import logging
import re
import time
from typing import Optional, List, Dict, Any, AsyncGenerator
from app.core.config import settings

import httpx

logger = logging.getLogger(__name__)

THINK_BLOCK_RE = re.compile(r"<think>.*?</think>", re.DOTALL)

# NVIDIA is the only supported AI provider for this project.
NVIDIA_MODELS: Dict[str, str] = {
    "chat": "nvidia/nemotron-3-super-120b-a12b",
    "stem": "nvidia/nemotron-3-super-120b-a12b",
    "coding": "nvidia/nemotron-3-super-120b-a12b",
    "vision": "nvidia/nemotron-3-super-120b-a12b",
    "embed": "nvidia/nv-embedqa-e5-v5",
}

# Cooldown before retrying a failed NVIDIA request.
PROVIDER_COOLDOWN = 30

# Shared persistent connection pool for all NVIDIA calls.
# Keeps TCP connections alive between requests — eliminates per-request
# TLS handshake overhead (~100-300 ms saved per call).
_HTTP_LIMITS = httpx.Limits(
    max_connections=50,
    max_keepalive_connections=20,
    keepalive_expiry=30,
)
_SHARED_HTTP_CLIENT: Optional[httpx.AsyncClient] = None

# Marker embedded in generated content. Older cached content (whitespace-stripped
# artifacts from nemotron-3.5 thinking mode) lacks it and is regenerated.
CACHE_VERSION = "<!--academia-v3-->"


def normalize_ai_markdown(text: str) -> str:
    """Clean up model output before caching.

    - Converts <center>mermaid ... </center> wrappers into proper fenced
      code blocks so diagrams actually render.
    - Strips stray HTML wrapper tags.
    - Inserts the cache-version marker.
    """
    if not text:
        return text
    text = re.sub(
        r"<center>\s*mermaid\s*([\s\S]*?)</center>",
        lambda m: "```mermaid\n" + m.group(1).strip() + "\n```",
        text,
        flags=re.IGNORECASE,
    )
    text = re.sub(r"</?center>", "", text)
    return CACHE_VERSION + "\n" + text.strip()


def _get_http_client() -> httpx.AsyncClient:
    global _SHARED_HTTP_CLIENT
    if _SHARED_HTTP_CLIENT is None or _SHARED_HTTP_CLIENT.is_closed:
        _SHARED_HTTP_CLIENT = httpx.AsyncClient(limits=_HTTP_LIMITS, timeout=120.0)
    return _SHARED_HTTP_CLIENT


def _strip_think(text: str) -> str:
    """Remove reasoning blocks some NVIDIA reasoning models may wrap in  thinking tags.

    IMPORTANT: must NOT strip() plain text — streamed chunks arrive one token
    at a time (often just ' word'), and stripping each chunk removes every
    leading space, concatenating the whole output into one unreadable blob.
    Only touch text that actually contains reasoning markers.
    """
    if not text:
        return text
    if "<think" not in text and not THINK_BLOCK_RE.search(text):
        return text
    stripped = THINK_BLOCK_RE.sub("", text)
    if "<think" in stripped:
        # Block was truncated (no closing tag) — keep whatever follows the opening tag
        idx = stripped.find("<think")
        stripped = stripped[idx + len("<think"):]
    return stripped.strip()


class EmptyResponseError(Exception):
    """Raised when a provider returns success but no usable content."""


def _make_client(api_key: Optional[str], base_url: str = None):
    if not api_key:
        return None
    import openai
    kwargs: Dict[str, Any] = {
        "api_key": api_key,
        # Reuse the shared persistent connection pool to avoid repeated
        # TLS handshakes on every request (~100-300 ms saving per call).
        "http_client": _get_http_client(),
    }
    if base_url:
        kwargs["base_url"] = base_url
    return openai.AsyncOpenAI(**kwargs)


def _normalize_content(content) -> str:
    """Workers AI may return content as a list of parts; flatten to text."""
    if content is None:
        return ""
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        if content and all(isinstance(p, dict) for p in content):
            if any("text" in p for p in content):
                parts = []
                for p in content:
                    text = p.get("text")
                    if text is not None:
                        parts.append(text if isinstance(text, str) else str(text))
                return "".join(parts)
            return json.dumps(content)
        return "".join(str(p) for p in content)
    return str(content)


class AIProvider:
    """NVIDIA-only AI provider wrapper.

    The project intentionally supports only NVIDIA. Any non-NVIDIA value is
    ignored and the runtime stays pinned to NVIDIA.
    """

    def __init__(self):
        self.provider = "nvidia"
        if settings.AI_PROVIDER and settings.AI_PROVIDER.lower() != "nvidia":
            logger.warning("Non-NVIDIA AI provider configured; forcing NVIDIA-only mode.")

        # Ordered list of usable providers: {name, client, models}
        self._providers: List[Dict[str, Any]] = []
        # Provider name -> timestamp until which it is skipped after a failure
        self._dead_providers: Dict[str, float] = {}

        self._init_providers()

    def _provider_config(self) -> Optional[Dict[str, Any]]:
        """Build the NVIDIA client + model table, or None if unconfigured."""
        if not (settings.NVIDIA_API_KEY or settings.NVIDIA_CODING_API_KEY
                or settings.NVIDIA_VISION_API_KEY):
            return None
        nvidia_base = settings.NVIDIA_BASE_URL
        if "integrate.api.nvidia.com" in nvidia_base and not nvidia_base.rstrip("/").endswith("/v1"):
            nvidia_base = nvidia_base.rstrip("/") + "/v1"
        client = _make_client(
            settings.NVIDIA_API_KEY
            or settings.NVIDIA_CODING_API_KEY
            or settings.NVIDIA_VISION_API_KEY,
            nvidia_base,
        )
        models = {
            "chat": settings.CHAT_MODEL,
            "stem": settings.STEM_MODEL,
            "coding": settings.CODING_MODEL,
            "vision": settings.VISION_MODEL,
            "embed": settings.EMBEDDING_MODEL,
        }
        return {"name": "nvidia", "client": client, "models": models}

    def _init_providers(self):
        """Initialize NVIDIA as the sole supported AI provider."""
        cfg = self._provider_config()
        self._providers = [cfg] if cfg else []
        if not self._providers:
            logger.warning("NVIDIA API key not configured. AI features will be unavailable.")

    def _live_providers(self) -> List[Dict[str, Any]]:
        now = time.time()
        return [
            p for p in self._providers
            if p["name"] not in self._dead_providers or self._dead_providers[p["name"]] <= now
        ]

    def _mark_dead(self, name: str):
        logger.error(f"AI provider '{name}' failed — cooling down for {PROVIDER_COOLDOWN}s")
        self._dead_providers[name] = time.time() + PROVIDER_COOLDOWN

    def _get_client(self, task: str = "chat"):
        for provider in self._providers:
            if provider["name"] == settings.AI_PROVIDER:
                return provider["client"]
        return self._providers[0]["client"] if self._providers else None

    def _get_model(self, task: str = "chat"):
        for provider in self._providers:
            if provider["name"] == settings.AI_PROVIDER:
                return provider["models"].get(task) or provider["models"]["chat"]
        if not self._providers:
            return settings.CHAT_MODEL
        return self._providers[0]["models"].get(task) or self._providers[0]["models"]["chat"]

    async def _chat_provider(
        self,
        provider: Dict[str, Any],
        messages: List[Dict[str, str]],
        temperature: float,
        max_tokens: int,
        stream: bool,
        task: str,
    ) -> AsyncGenerator[str, None]:
        client = provider["client"]
        model = provider["models"].get(task) or provider["models"]["chat"]

        extra_kwargs: Dict[str, Any] = {}
        if "nemotron-3-super" in model:
            # Nemotron-3-super emits high-quality reasoning output; a tight
            # reasoning budget keeps time-to-first-token low.
            thinking_budget = 512 if task == "coding" else 384
            extra_kwargs = {
                "extra_body": {
                    "chat_template_kwargs": {"enable_thinking": True},
                    "reasoning_budget": thinking_budget,
                }
            }
        elif "nemotron-3.5" in model or "lightning" in model:
            # Nemotron-3.5 Lightning DROPS ALL WHITESPACE from streamed output
            # when thinking is enabled (verified: 0 spaces / 0 newlines), which
            # corrupts markdown. Thinking off fixes formatting and is ~5x faster.
            extra_kwargs = {
                "extra_body": {
                    "chat_template_kwargs": {"enable_thinking": False}
                }
            }

        try:
            response = await client.chat.completions.create(
                model=model,
                messages=messages,
                temperature=temperature,
                max_tokens=max_tokens,
                stream=stream,
                **extra_kwargs,
            )
        except Exception as e:
            # NVIDIA returns 404 when the base URL lacks the /v1 suffix
            # (or points at a stale endpoint). Rebuild the client against the
            # canonical endpoint and retry once.
            if provider["name"] == "nvidia" and "404" in str(e):
                import openai as _openai
                canonical = "https://integrate.api.nvidia.com/v1"
                client = _openai.AsyncOpenAI(api_key=client.api_key, base_url=canonical)
                response = await client.chat.completions.create(
                    model=model,
                    messages=messages,
                    temperature=temperature,
                    max_tokens=max_tokens,
                    stream=stream,
                    **extra_kwargs,
                )
            else:
                raise
        if stream:
            saw_content = False
            async for chunk in response:
                if not chunk.choices:
                    continue
                content = chunk.choices[0].delta.content
                if content:
                    saw_content = True
                    yield json.dumps({
                        "role": "assistant",
                        "content": _strip_think(_normalize_content(content)),
                    })
            if not saw_content:
                raise EmptyResponseError(
                    f"stream ended without content (model {model})"
                )
        else:
            if not response.choices:
                raise EmptyResponseError(
                    f"empty response (model {model})"
                )
            content = _strip_think(_normalize_content(response.choices[0].message.content))
            if not content:
                raise EmptyResponseError(
                    f"empty response (model {model}, finish_reason={response.choices[0].finish_reason})"
                )
            yield json.dumps({"role": "assistant", "content": content})

    async def chat(
        self,
        messages: List[Dict[str, str]],
        temperature: float = 0.7,
        max_tokens: int = 2048,
        stream: bool = True,
        task: str = "chat",
    ) -> AsyncGenerator[str, None]:
        providers = self._live_providers()
        if not providers:
            if self._providers:
                # Provider exists but is in cooldown — retry now (may have recovered)
                providers = [self._providers[0]]
            else:
                yield json.dumps({
                    "role": "assistant",
                    "content": "I'm running in offline mode. Please configure an NVIDIA API key."
                })
                return

        errors = []
        for provider in providers:
            try:
                async for chunk in self._chat_provider(
                    provider, messages, temperature, max_tokens, stream, task
                ):
                    yield chunk
                return  # completed without error
            except EmptyResponseError as e:
                # Reasoning model returned empty — retry once with a larger budget
                logger.warning(
                    f"AI empty response ({provider['name']} / "
                    f"{provider['models'].get(task, 'chat')}, budget={max_tokens}): {e} "
                    f"— retrying with 4x budget"
                )
                try:
                    async for chunk in self._chat_provider(
                        provider, messages, temperature, max_tokens * 4, stream, task
                    ):
                        yield chunk
                    return
                except Exception as e2:
                    errors.append(f"{provider['name']}: {e2}")
                    logger.error(f"AI chat retry failed: {e2}")
                    self._mark_dead(provider["name"])
            except Exception as e:
                errors.append(f"{provider['name']}: {e}")
                logger.error(f"AI chat error ({provider['name']} / "
                             f"{provider['models'].get(task, 'chat')}): {e}")
                self._mark_dead(provider["name"])

        yield json.dumps({
            "role": "assistant",
            "content": f"AI request failed: {'; '.join(errors)}"
        })

    async def generate_embeddings(self, texts: List[str]) -> List[List[float]]:
        dim = settings.EMBEDDING_DIMENSION
        if getattr(self, "_embeddings_failed", False):
            return [[0.0] * dim for _ in texts]
        import asyncio as _asyncio

        for provider in self._live_providers():
            embed_model = provider["models"].get("embed")
            client = provider["client"]
            if not embed_model or client is None:
                continue
            try:
                batches = [texts[i:i + 128] for i in range(0, len(texts), 128)]
                kwargs = {}
                if "e5" in embed_model:
                    # NVIDIA asymmetric embed models require input_type
                    kwargs["extra_body"] = {"input_type": "passage"}
                results = await _asyncio.gather(*[
                    client.embeddings.create(model=embed_model, input=batch, **kwargs)
                    for batch in batches
                ])
                vectors = []
                for res in results:
                    vectors.extend(item.embedding for item in res.data)
                return vectors
            except Exception as e:
                logger.error(f"Embedding error ({provider['name']} / {embed_model}): {e}")
                self._mark_dead(provider["name"])
        # Every configured provider failed — short-circuit all future calls
        self._embeddings_failed = True
        return [[0.0] * dim for _ in texts]

    async def extract_text_from_file(self, file_path: str, file_type: str) -> str:
        text = ""
        try:
            if file_type == "pdf":
                import pdfplumber
                with pdfplumber.open(file_path) as pdf:
                    text = "\n".join(page.extract_text() or "" for page in pdf.pages)
            elif file_type in ("docx", "doc"):
                from docx import Document
                doc = Document(file_path)
                text = "\n".join(p.text for p in doc.paragraphs)
            elif file_type in ("pptx", "ppt"):
                from pptx import Presentation
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            text += shape.text + "\n"
            elif file_type in ("xlsx", "xls"):
                import openpyxl
                wb = openpyxl.load_workbook(file_path, data_only=True)
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    for row in ws.iter_rows(values_only=True):
                        text += " | ".join(str(cell) if cell else "" for cell in row) + "\n"
            elif file_type == "csv":
                import csv
                with open(file_path, "r", encoding="utf-8") as f:
                    reader = csv.reader(f)
                    text = "\n".join(" | ".join(row) for row in reader)
            elif file_type in ("txt", "md"):
                with open(file_path, "r", encoding="utf-8") as f:
                    text = f.read()
            elif file_type in ("png", "jpg", "jpeg", "gif", "bmp", "webp"):
                text = "[Image file — no text content extracted]"
            elif file_type in ("mp4", "avi", "mov", "mkv", "webm"):
                text = "[Video file — no text content extracted]"
            elif file_type in ("mp3", "wav", "ogg", "m4a", "flac"):
                text = "[Audio file — no text content extracted]"
            elif file_type == "html":
                from bs4 import BeautifulSoup
                with open(file_path, "r", encoding="utf-8") as f:
                    soup = BeautifulSoup(f.read(), "lxml")
                    text = soup.get_text(separator="\n")
            elif file_type == "epub":
                import ebooklib
                from ebooklib import epub
                from bs4 import BeautifulSoup
                book = epub.read_epub(file_path)
                for item in book.get_items():
                    if item.get_type() == ebooklib.ITEM_DOCUMENT:
                        soup = BeautifulSoup(item.get_content(), "lxml")
                        text += soup.get_text() + "\n"
            else:
                text = "Unsupported file type for text extraction."
        except Exception as e:
            logger.error(f"Text extraction error for {file_type}: {e}")
            text = f"Error extracting text: {str(e)}"
        return text

    async def generate_summary(self, content: str, summary_type: str = "short_notes", max_length: int = 500) -> Dict[str, Any]:
        prompts = {
            "short_notes": "Generate concise short notes from the following content. Include key points, definitions, and important concepts.",
            "cheat_sheet": "Create a one-page cheat sheet covering the essential formulas, concepts, and key points.",
            "flashcards": "Generate flashcards in Q&A format from the content.",
            "mind_map": "Create a hierarchical mind map structure from the content.",
            "formula_sheet": "Extract all formulas, equations, and mathematical expressions.",
            "revision_notes": "Create comprehensive revision notes organized by topics and subtopics.",
            "important_questions": "Generate a list of important exam questions based on the content.",
            "one_page_summary": "Summarize the entire content into a single page.",
        }
        prompt = prompts.get(summary_type, prompts["short_notes"])
        messages = [
            {"role": "system", "content": f"You are an expert educator and summarizer. {prompt} Format the output in markdown. Keep it under {max_length} words."},
            {"role": "user", "content": f"Content to summarize:\n\n{content[:10000]}"},
        ]
        result = ""
        # qwen reasoning models spend tokens on thinking — give a generous budget and retry once if truncated
        for _ in range(2):
            result = ""
            async for chunk in self.chat(messages, temperature=0.3, max_tokens=max(2048, max_length * 4), task="chat"):
                data = json.loads(chunk)
                result += data.get("content", "")
            if result.strip():
                break
        return {
            "title": summary_type.replace("_", " ").title(),
            "content": result,
            "summary_type": summary_type,
            "format": "markdown",
        }

    def _fallback_questions(
        self, content: str, question_type: str = "mcq", count: int = 10, difficulty: str = "medium"
    ) -> List[Dict[str, Any]]:
        text = re.sub(r"\s+", " ", content or "").strip()
        snippet = text[:800] if text else "core concepts and study methods"
        tokens = re.findall(r"[A-Za-z][A-Za-z0-9]{2,}", snippet)
        focus_terms = [term for term in tokens if len(term) > 3][:6] or ["concept", "method", "practice", "review", "understanding", "application"]

        questions: List[Dict[str, Any]] = []
        for i in range(max(1, min(count, 10))):
            topic = focus_terms[i % len(focus_terms)]
            correct = f"{topic.title()} is a central idea in this topic"
            distractors = [
                f"{topic.title()} is unrelated to the lesson",
                "Only memorizing words is enough",
                "Skipping practice has no effect on understanding",
            ]
            if question_type.lower() == "mcq":
                options = [correct, *distractors]
                # Keep a valid answer list even when the AI output is weird.
                if len(options) < 4:
                    options += ["Reviewing examples helps build confidence"]
                correct_answer = correct
            else:
                options = ["Short answer response"]
                correct_answer = "Apply the key principle from the lesson"

            questions.append({
                "question_text": f"Which statement best reflects the main idea of this topic in relation to '{topic}'?",
                "question_type": question_type or "mcq",
                "options": options,
                "correct_answer": correct_answer,
                "explanation": f"The lesson emphasizes {topic} as a key concept, and the strongest answer matches the core idea rather than a distractor.",
                "difficulty": 2 if difficulty == "easy" else 3 if difficulty == "medium" else 4,
                "marks": 1,
            })
        return questions[:count]

    async def generate_questions(
        self, content: str, question_type: str = "mcq", count: int = 10, difficulty: str = "medium"
    ) -> List[Dict[str, Any]]:
        prompt = f"""
Generate {count} {question_type} questions at {difficulty} difficulty based on the content.
For each question provide:
- question_text
- question_type
- options (array of choices for MCQ)
- correct_answer
- explanation
- difficulty (1-5)
- marks

Format as JSON array.
"""
        messages = [
            {"role": "system", "content": "You are an expert exam question generator. Always respond with valid JSON only."},
            {"role": "user", "content": f"{prompt}\n\nContent:\n{content[:8000]}"},
        ]
        result = ""
        # Use STEM model for question generation (stronger reasoning)
        try:
            async for chunk in self.chat(messages, temperature=0.4, max_tokens=max(2048, count * 512), task="stem"):
                data = json.loads(chunk)
                result += data.get("content", "")
        except Exception as e:
            logger.warning(f"Question generation chat failed; using fallback questions: {e}")
            return self._fallback_questions(content, question_type, count, difficulty)

        try:
            if "```json" in result:
                result = result.split("```json")[1].split("```")[0].strip()
            elif "```" in result:
                result = result.split("```")[1].split("```")[0].strip()
            questions = json.loads(result)
            if isinstance(questions, list) and questions:
                return questions[:count]
            if isinstance(questions, dict) and questions:
                return [questions]
        except (TypeError, ValueError, json.JSONDecodeError):
            pass

        try:
            start, end = result.index("["), result.rindex("]")
            questions = json.loads(result[start:end + 1])
            if isinstance(questions, list) and questions:
                return questions[:count]
        except (ValueError, json.JSONDecodeError):
            pass

        logger.warning(f"Failed to parse questions from AI output; using fallback questions. Output preview: {result[:200]}")
        return self._fallback_questions(content, question_type, count, difficulty)


ai_provider = AIProvider()